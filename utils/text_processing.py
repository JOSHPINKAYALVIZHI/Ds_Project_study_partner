import re
from pptx import Presentation
from PyPDF2 import PdfReader

def extract_text_from_pdf(file):
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return clean_text(text)

def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return clean_text(text)

def clean_text(text):
    lines = text.split("\n")
    clean_lines = []
    for line in lines:
        line = line.strip()
        if not line:
            continue
        # Remove headers/footers 
        if line.isupper():
            continue
        if re.fullmatch(r'\d+', line):  # page numbers
            continue
        if any(k in line.upper() for k in ["LECTURE", "SLIDE", "MARKS", "INSTITUTE", "DATE", "NAME"]):
            continue
        # Remove special chars
        line = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', line)
        clean_lines.append(line)
    return " ".join(clean_lines)
